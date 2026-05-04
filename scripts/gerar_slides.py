"""
Gera apresentação executiva em PowerPoint (storytelling, sem código).
Uso: python scripts/gerar_slides.py
Saída: reports/Apresentacao_NPS_Fase1.pptx
"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "reports" / "Apresentacao_NPS_Fase1.pptx"


def add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle


def add_bullets(prs: Presentation, title: str, bullets: list[str]) -> None:
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    for i, text in enumerate(bullets):
        p = body.paragraphs[0] if i == 0 else body.add_paragraph()
        p.text = text
        p.level = 0
        p.font.size = Pt(20)


def main() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_title_slide(
        prs,
        "NPS Preditivo — Tech Challenge Fase 1",
        "E-commerce: da operação à experiência do cliente\nFIAP · Storytelling para gestão",
    )

    add_bullets(
        prs,
        "Contexto do problema",
        [
            "E-commerce em crescimento: mais pedidos, entregas e interações.",
            "NPS varia muito entre clientes — e a pesquisa só vem no fim da jornada.",
            "Consequência: a empresa reage tarde; queremos antecipar problemas.",
        ],
    )

    add_bullets(
        prs,
        "Pergunta que orienta a análise",
        [
            "Quais fatores operacionais mais explicam o NPS?",
            "Onde estão os detratores e qual o “ponto de ruptura” da experiência?",
            "(Complemento) Qual pedido tem maior risco de NPS baixo antes da pesquisa?",
        ],
    )

    add_bullets(
        prs,
        "Principais insights da EDA",
        [
            "Base com ~2.500 pedidos: NPS médio baixo; predominam detrators na escala 0–10.",
            "Atraso na entrega: sem atraso, nota média muito maior; com 4+ dias, nota despenca.",
            "Reclamações e contatos com atendimento acompanham queda forte do NPS.",
            "Regiões: diferenças menores — o gargalo é execução, não “só marketing regional”.",
        ],
    )

    add_bullets(
        prs,
        "Fatores que mais impactam a satisfação",
        [
            "Logística: dias de atraso e tentativas de entrega.",
            "Qualidade / volume de reclamações e tempo de resolução.",
            "Pressão no SAC: cada contato extra sinaliza jornada problemática.",
            "CSAT interno acompanha o NPS — bom alinhamento de indicadores.",
        ],
    )

    add_bullets(
        prs,
        "Recomendações práticas",
        [
            "Plano de ação para atraso ≥ 4 dias: comunicação proativa + compensação.",
            "Priorizar root-cause de reclamações (produto, embalagem, transportadora, promessa).",
            "Reduzir necessidade de contato: rastreio claro e expectativa de prazo.",
            "Pilotar fila de risco com score preditivo antes da pesquisa NPS.",
        ],
    )

    add_bullets(
        prs,
        "Modelo preditivo (visão executiva)",
        [
            "Usa dados de pedido, logística e atendimento para estimar o NPS.",
            "Ajuda a priorizar onde logística e SAC agem antes da pesquisa.",
            "Variáveis mais relevantes: atraso, reclamações e sinais de fricção no atendimento.",
        ],
    )

    add_bullets(
        prs,
        "Limitações e riscos",
        [
            "Correlação não implica causa: validar com pilotos e experimentos.",
            "Modelo pode mudar com sazonalidade, greves ou mudança de transportadora.",
            "Uso responsável: apoio à decisão, não punição automática de times ou clientes.",
        ],
    )

    add_bullets(
        prs,
        "Próximos passos",
        [
            "Inserir nos slides as figuras da pasta reports/ (gráficos da EDA e do modelo).",
            "Definir piloto em uma região ou categoria e acompanhar NPS real vs predito.",
            "Alinhar OKRs de logística e SAC às alavancas identificadas.",
        ],
    )

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"Slides gerados: {OUT}")


if __name__ == "__main__":
    main()
